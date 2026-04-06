from __future__ import annotations

import csv
import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable

from langchain_core.documents import Document


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        stripped = data.strip()
        if stripped:
            self.parts.append(stripped)

    def get_text(self) -> str:
        return "\n".join(self.parts)


def load_text_document(path: str | Path, *, encoding: str = "utf-8") -> list[Document]:
    file_path = Path(path)
    return [
        Document(
            page_content=file_path.read_text(encoding=encoding),
            metadata={"source": str(file_path)},
        )
    ]


def load_csv_documents(path: str | Path) -> list[Document]:
    file_path = Path(path)
    with file_path.open("r", encoding="utf-8", newline="") as file:
        reader = csv.DictReader(file)
        return [
            Document(
                page_content="\n".join(f"{key}: {value}" for key, value in row.items()),
                metadata={"source": str(file_path), "row": index},
            )
            for index, row in enumerate(reader, start=1)
        ]


def load_json_documents(path: str | Path) -> list[Document]:
    file_path = Path(path)
    data = json.loads(file_path.read_text(encoding="utf-8"))

    if isinstance(data, list):
        items = data
    else:
        items = [data]

    documents: list[Document] = []
    for index, item in enumerate(items, start=1):
        documents.append(
            Document(
                page_content=json.dumps(item, ensure_ascii=False, indent=2),
                metadata={"source": str(file_path), "index": index},
            )
        )
    return documents


def load_html_document(path: str | Path) -> list[Document]:
    file_path = Path(path)
    parser = _HTMLTextExtractor()
    parser.feed(file_path.read_text(encoding="utf-8"))
    return [
        Document(
            page_content=parser.get_text(),
            metadata={"source": str(file_path)},
        )
    ]


def load_pdf_documents(path: str | Path) -> list[Document]:
    from pypdf import PdfReader

    file_path = Path(path)
    reader = PdfReader(str(file_path))

    documents: list[Document] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        documents.append(
            Document(
                page_content=text.strip(),
                metadata={"source": str(file_path), "page": page_number},
            )
        )
    return documents


def load_excel_documents(path: str | Path) -> list[Document]:
    import pandas as pd

    file_path = Path(path)
    sheets = pd.read_excel(file_path, sheet_name=None)

    documents: list[Document] = []
    for sheet_name, dataframe in sheets.items():
        for row_index, row in dataframe.fillna("").iterrows():
            documents.append(
                Document(
                    page_content="\n".join(
                        f"{column}: {value}" for column, value in row.items()
                    ),
                    metadata={
                        "source": str(file_path),
                        "sheet": sheet_name,
                        "row": int(row_index) + 1,
                    },
                )
            )
    return documents


def load_word_documents(path: str | Path) -> list[Document]:
    import docx2txt

    file_path = Path(path)
    return [
        Document(
            page_content=docx2txt.process(str(file_path)).strip(),
            metadata={"source": str(file_path)},
        )
    ]


def load_pptx_documents(path: str | Path) -> list[Document]:
    from pptx import Presentation

    file_path = Path(path)
    presentation = Presentation(str(file_path))
    documents: list[Document] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                texts.append(text)

        documents.append(
            Document(
                page_content="\n".join(texts),
                metadata={"source": str(file_path), "slide": slide_number},
            )
        )
    return documents


def load_arxiv_documents(query: str, *, load_max_docs: int = 2) -> list[Document]:
    import arxiv

    search = arxiv.Search(
        query=query,
        max_results=load_max_docs,
        sort_by=arxiv.SortCriterion.Relevance,
    )

    documents: list[Document] = []
    client = arxiv.Client()
    for index, result in enumerate(client.results(search), start=1):
        documents.append(
            Document(
                page_content=result.summary.strip(),
                metadata={
                    "source": result.entry_id,
                    "title": result.title,
                    "published": result.published.isoformat(),
                    "authors": [author.name for author in result.authors],
                    "rank": index,
                },
            )
        )
    return documents


def load_directory_documents(directory: str | Path) -> list[Document]:
    base_dir = Path(directory)
    documents: list[Document] = []

    for path in sorted(base_dir.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix == ".txt":
            documents.extend(load_text_document(path))
        elif path.suffix == ".csv":
            documents.extend(load_csv_documents(path))
        elif path.suffix == ".json":
            documents.extend(load_json_documents(path))
        elif path.suffix in {".html", ".htm"}:
            documents.extend(load_html_document(path))
        elif path.suffix == ".pdf":
            documents.extend(load_pdf_documents(path))
        elif path.suffix in {".xlsx", ".xls"}:
            documents.extend(load_excel_documents(path))
        elif path.suffix == ".docx":
            documents.extend(load_word_documents(path))
        elif path.suffix == ".pptx":
            documents.extend(load_pptx_documents(path))

    return documents


__all__ = [
    "load_arxiv_documents",
    "load_csv_documents",
    "load_directory_documents",
    "load_excel_documents",
    "load_html_document",
    "load_json_documents",
    "load_pdf_documents",
    "load_pptx_documents",
    "load_text_document",
    "load_word_documents",
]
